from __future__ import annotations

import logging
import re
import shutil
import sys
from pathlib import Path
from collections.abc import Callable

from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader
from langchain_core.documents import Document
from langchain_text_splitters import (
    MarkdownHeaderTextSplitter,
    RecursiveCharacterTextSplitter,
)

from uacragent.domain.errors import IngestError
from uacragent.domain.types import DocumentType
from uacragent.infra.settings import Settings
from uacragent.infra.workspace import WorkspacePaths

logger = logging.getLogger(__name__)

# Maximum iterations for filename collision resolution.  At 1 000 a collision
# that cannot be resolved is practically impossible with human-uploaded files;
# anything beyond suggests a filesystem problem rather than a collision.
_MAX_COLLISION_COUNTER = 1_000

# Pre-load file-size thresholds.
# Files larger than _WARN_BYTES are logged as warnings; they may be slow.
# Files larger than _BLOCK_BYTES are refused entirely — loading them risks OOM.
_LARGE_FILE_WARN_BYTES  = 100 * 1024 * 1024   # 100 MB
_LARGE_FILE_BLOCK_BYTES = 300 * 1024 * 1024   # 300 MB


# ---------------------------------------------------------------------------
# Multi-stage splitting pipeline definitions
# ---------------------------------------------------------------------------
# Each document type has a pipeline: an ordered list of splitting stages.
# A stage is a callable  list[Document] -> list[Document].
#
# Strategy rationale:
#   Textbook   – header split to isolate chapters/sections, then recursive
#                character split to break long sections into retrievable chunks.
#   Syllabus   – header split to isolate policy/schedule sections, then small
#                character split so each policy item is its own chunk.
#   Lecture    – sentence-aware recursive split (uses sentence separators first)
#                to keep slide-level ideas intact.
#   Past Exam  – regex split on question boundaries (numbered items) to keep
#                each question together, then small character split for long Qs.
#   Assignment – same boundary approach as past exams but tuned for problem
#                sets (Problem / Question / Exercise headers).
#   Other      – standard recursive character split as a safe default.
# ---------------------------------------------------------------------------


def _make_header_then_recursive(
    headers_to_split_on: list[tuple[str, str]],
    chunk_size: int,
    chunk_overlap: int,
) -> list[Callable[[list[Document]], list[Document]]]:
    """Return a 2-stage pipeline: markdown-header split, then recursive."""

    def _header_stage(docs: list[Document]) -> list[Document]:
        md_splitter = MarkdownHeaderTextSplitter(
            headers_to_split_on=headers_to_split_on,
            strip_headers=False,
        )
        out: list[Document] = []
        for doc in docs:
            splits = md_splitter.split_text(doc.page_content)
            for split in splits:
                # Merge parent metadata into each split
                merged_meta = {**doc.metadata, **split.metadata}
                out.append(Document(page_content=split.page_content, metadata=merged_meta))
        return out

    def _recursive_stage(docs: list[Document]) -> list[Document]:
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
        )
        return splitter.split_documents(docs)

    return [_header_stage, _recursive_stage]


def _make_sentence_aware_recursive(
    chunk_size: int,
    chunk_overlap: int,
) -> list[Callable[[list[Document]], list[Document]]]:
    """Single-stage pipeline with sentence-aware separators for lecture notes."""

    def _split(docs: list[Document]) -> list[Document]:
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=[
                "\n\n",        # paragraph break
                "\n",          # line break
                ". ",          # sentence end
                "? ",          # question end
                "! ",          # exclamation end
                "; ",          # semicolon (list items, definitions)
                ", ",          # comma
                " ",           # word boundary
                "",            # character fallback
            ],
        )
        return splitter.split_documents(docs)

    return [_split]


def _make_question_boundary_split(
    boundary_pattern: str,
    chunk_size: int,
    chunk_overlap: int,
) -> list[Callable[[list[Document]], list[Document]]]:
    """2-stage pipeline: regex split on question/problem boundaries, then
    recursive character split for any chunks that are still too large."""

    compiled = re.compile(boundary_pattern)

    def _boundary_stage(docs: list[Document]) -> list[Document]:
        out: list[Document] = []
        for doc in docs:
            parts = compiled.split(doc.page_content)
            for part in parts:
                part = part.strip()
                if part:
                    out.append(Document(page_content=part, metadata=dict(doc.metadata)))
        return out

    def _recursive_stage(docs: list[Document]) -> list[Document]:
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
        )
        return splitter.split_documents(docs)

    return [_boundary_stage, _recursive_stage]


def _make_recursive_only(
    chunk_size: int,
    chunk_overlap: int,
) -> list[Callable[[list[Document]], list[Document]]]:
    """Single-stage standard recursive split."""

    def _split(docs: list[Document]) -> list[Document]:
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
        )
        return splitter.split_documents(docs)

    return [_split]


# ---------------------------------------------------------------------------
# Pipeline registry – one pipeline per document type
# ---------------------------------------------------------------------------

def _get_pipeline(doc_type: DocumentType) -> list[Callable[[list[Document]], list[Document]]]:
    """Return the multi-stage splitting pipeline for the given document type."""

    if doc_type == DocumentType.textbook:
        # Stage 1: split on markdown headers (chapter / section / subsection)
        # Stage 2: recursive split to max 1500-char retrievable chunks
        return _make_header_then_recursive(
            headers_to_split_on=[
                ("#", "chapter"),
                ("##", "section"),
                ("###", "subsection"),
            ],
            chunk_size=1500,
            chunk_overlap=200,
        )

    if doc_type == DocumentType.syllabus:
        # Stage 1: header split to isolate policy / schedule / grading sections
        # Stage 2: small recursive split so each rule/policy is a chunk
        return _make_header_then_recursive(
            headers_to_split_on=[
                ("#", "section"),
                ("##", "subsection"),
            ],
            chunk_size=800,
            chunk_overlap=100,
        )

    if doc_type == DocumentType.lecture_note:
        # Sentence-aware split keeps slide bullet points together
        return _make_sentence_aware_recursive(
            chunk_size=1000,
            chunk_overlap=150,
        )

    if doc_type == DocumentType.past_exam:
        # Stage 1: split on question boundaries  (1. / Q1 / Question 1 / (a) etc.)
        # Stage 2: recursive split for any oversized questions
        return _make_question_boundary_split(
            boundary_pattern=(
                r"(?=\n\s*(?:"
                r"\d+[\.\)]\s"          # "1. " or "1) "
                r"|[Qq](?:uestion)?\s*\d+"  # "Q1", "Question 1"
                r"|Part\s+[A-Za-z0-9]"  # "Part A", "Part 1"
                r"|\([a-z]\)"           # "(a)"
                r"))"
            ),
            chunk_size=500,
            chunk_overlap=80,
        )

    if doc_type == DocumentType.assignment:
        # Stage 1: split on problem/exercise boundaries
        # Stage 2: recursive split for long problems
        return _make_question_boundary_split(
            boundary_pattern=(
                r"(?=\n\s*(?:"
                r"(?:Problem|Exercise|Task|Question)\s*\d+"  # "Problem 1"
                r"|\d+[\.\)]\s"        # "1. " or "1) "
                r"|\([a-z]\)"          # "(a)"
                r"))"
            ),
            chunk_size=600,
            chunk_overlap=100,
        )

    # DocumentType.other — safe default
    return _make_recursive_only(chunk_size=1000, chunk_overlap=150)


def run_splitting_pipeline(
    docs: list[Document],
    doc_type: DocumentType,
) -> list[Document]:
    """Execute the multi-stage splitting pipeline for the given document type.

    Each stage receives the output of the previous stage.  Metadata for
    ``doc_type`` is stamped on every resulting chunk.
    """
    pipeline = _get_pipeline(doc_type)
    current = docs
    for stage in pipeline:
        current = stage(current)

    # Stamp doc_type on every chunk
    for chunk in current:
        chunk.metadata["doc_type"] = doc_type.value

    return current


# ---------------------------------------------------------------------------
# OCR helpers (used by _load_pptx)
# ---------------------------------------------------------------------------

def _check_ocr_available() -> bool:
    """Return True when both pytesseract and PIL are importable AND Tesseract
    is reachable (either bundled via rthook_tesseract.py or system-installed).

    In a frozen PyInstaller build the Tesseract binary is bundled in _MEIPASS
    and rthook_tesseract.py sets pytesseract.tesseract_cmd before this code
    runs, so the get_tesseract_version() call will succeed.
    """
    try:
        import pytesseract as _pyt  # noqa: PLC0415
        from PIL import Image as _PILImage  # noqa: PLC0415, F401
        _pyt.get_tesseract_version()
        return True
    except Exception:
        return False


# python-pptx MSO_SHAPE_TYPE.PICTURE value.  Defined as a module constant so
# the magic number appears exactly once and is clearly labelled.
_MSO_PICTURE_TYPE: int = 13


def _ocr_slide_images(slide: object) -> list[str]:
    """Run Tesseract OCR on all image shapes in *slide*.

    Returns a list of non-empty OCR text strings (one per image that
    yielded readable text).  Returns an empty list on any failure so
    callers can degrade to text-only without crashing.
    """
    results: list[str] = []
    try:
        import pytesseract as _pyt  # noqa: PLC0415
        from PIL import Image as _PILImage  # noqa: PLC0415
        import io as _io

        for shape in slide.shapes:  # type: ignore[attr-defined]
            try:
                # Only process picture shapes that expose image bytes.
                if shape.shape_type != _MSO_PICTURE_TYPE:
                    continue
                img_bytes = shape.image.blob
                img = _PILImage.open(_io.BytesIO(img_bytes)).convert("RGB")
                text = _pyt.image_to_string(img, lang="eng").strip()
                if text:
                    results.append(text)
            except Exception:  # noqa: BLE001
                continue
    except Exception:  # noqa: BLE001
        pass
    return results


# ---------------------------------------------------------------------------
# DocumentLoader
# ---------------------------------------------------------------------------

class DocumentLoader:
    """Loads and processes documents with multi-stage type-aware splitting."""

    def __init__(self, settings: Settings) -> None:
        self.settings = settings

    def load_single_file(self, path: str) -> list[Document]:
        """Load a single file and return raw LangChain Documents.

        Supported formats
        -----------------
        - PDF   (.pdf)             -- PyPDFLoader (text-layer extraction)
        - Word  (.docx)            -- Docx2txtLoader
        - Slides(.pptx, .ppt)      -- python-pptx + pytesseract OCR
        - Text  (.txt, .md, .py, .js, .ts, .html, .htm, .xml, .json)
        - CSV   (.csv)             -- row-chunked plain-text table
        """
        try:
            ext = Path(path).suffix.lower()
            if ext == ".pdf":
                return PyPDFLoader(path).load()
            elif ext in {".txt", ".md", ".py", ".js", ".ts",
                         ".html", ".htm", ".xml", ".json"}:
                # Read directly with errors="replace" so non-UTF-8 files
                # (Latin-1, GBK, etc.) are loaded with replacement characters
                # instead of raising UnicodeDecodeError.
                content = Path(path).read_text(encoding="utf-8", errors="replace")
                return [Document(page_content=content, metadata={"source": path})]
            elif ext == ".docx":
                return Docx2txtLoader(path).load()
            elif ext == ".csv":
                return self._load_csv(path)
            elif ext in {".pptx", ".ppt"}:
                return self._load_pptx(path)
            else:
                raise IngestError(f"Unsupported file type: {path}")
        except IngestError:
            raise
        except Exception as exc:  # noqa: BLE001
            raise IngestError(f"Failed to load {path}: {exc}") from exc

    @staticmethod
    def _load_pptx(path: str) -> list[Document]:
        """Load a PowerPoint file (.pptx or .ppt) as a list of Documents.

        Extraction strategy
        -------------------
        For each slide a single Document is produced containing:
          - Slide number header
          - Title text (if present)
          - All text-frame body content (bullets, labels, notes boxes)
          - Speaker notes
          - OCR text from any image shapes (using pytesseract + PIL)

        OCR is attempted when:
          - python-pptx is available (required)
          - pytesseract AND Pillow are available (optional -- degrades to
            text-only when absent or when Tesseract is not installed)

        Legacy binary .ppt files are NOT supported by python-pptx; an
        IngestError with a clear conversion message is raised instead.
        """
        ext = Path(path).suffix.lower()

        # ------------------------------------------------------------------
        # .ppt binary format: python-pptx cannot open it.  Give a clear,
        # actionable error rather than a confusing low-level exception.
        # ------------------------------------------------------------------
        if ext == ".ppt":
            raise IngestError(
                f"'{Path(path).name}' is in the legacy PowerPoint 97-2003 "
                "binary format (.ppt) which cannot be read directly.\n"
                "Please open the file in PowerPoint or LibreOffice Impress "
                "and save it as a modern .pptx file, then re-add it."
            )

        # ------------------------------------------------------------------
        # python-pptx is a required dependency for .pptx loading.
        # ------------------------------------------------------------------
        try:
            from pptx import Presentation  # type: ignore[import-untyped]
            from pptx.util import Pt as _Pt   # noqa: F401 -- validates import
        except ImportError as exc:
            raise IngestError(
                "python-pptx is required to load .pptx files. "
                "Install it with:  pip install python-pptx"
            ) from exc

        # ------------------------------------------------------------------
        # Determine whether OCR is available.
        # ------------------------------------------------------------------
        _ocr_available = _check_ocr_available()

        # ------------------------------------------------------------------
        # Open and iterate slides.
        # ------------------------------------------------------------------
        try:
            prs = Presentation(path)
        except Exception as exc:  # noqa: BLE001
            raise IngestError(
                f"Could not open '{Path(path).name}': {exc}. "
                "The file may be corrupt or password-protected."
            ) from exc

        docs: list[Document] = []
        fname = Path(path).name

        for slide_idx, slide in enumerate(prs.slides, start=1):
            parts: list[str] = [f"[Slide {slide_idx}]"]

            # -- Title ------------------------------------------------
            try:
                if slide.shapes.title and slide.shapes.title.has_text_frame:
                    title_text = slide.shapes.title.text_frame.text.strip()
                    if title_text:
                        parts.append(f"Title: {title_text}")
            except Exception:  # noqa: BLE001
                pass

            # -- Text frames (body, labels, text boxes) ---------------
            try:
                for shape in slide.shapes:
                    # Skip the title shape -- already handled above.
                    try:
                        if shape is slide.shapes.title:
                            continue
                    except Exception:  # noqa: BLE001
                        pass
                    if shape.has_text_frame:
                        text = shape.text_frame.text.strip()
                        if text:
                            parts.append(text)
            except Exception:  # noqa: BLE001
                pass

            # -- Speaker notes ----------------------------------------
            try:
                if (
                    slide.has_notes_slide
                    and slide.notes_slide.notes_text_frame
                ):
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append(f"[Notes] {notes}")
            except Exception:  # noqa: BLE001
                pass

            # -- Image OCR (if available) / image-count note (if not) -----
            if _ocr_available:
                ocr_texts = _ocr_slide_images(slide)
                if ocr_texts:
                    parts.append("[Image text (OCR)]")
                    parts.extend(ocr_texts)
            else:
                # OCR not available — count embedded picture shapes so the
                # indexed text at least records that visual content exists.
                # This lets the LLM answer questions like "how many images are
                # on slide 3?" and avoids silently empty slides in image-heavy
                # presentations.
                try:
                    img_count = sum(
                        1 for shape in slide.shapes
                        if shape.shape_type == _MSO_PICTURE_TYPE
                    )
                    if img_count > 0:
                        parts.append(
                            f"[{img_count} image"
                            f"{'s' if img_count != 1 else ''} on this slide"
                            f" — text inside images not extracted]"
                        )
                except Exception:  # noqa: BLE001
                    pass

            content = "\n".join(parts)
            if content.strip():
                docs.append(Document(
                    page_content=content,
                    metadata={
                        "source":      path,
                        "slide_number": slide_idx,
                        "filename":    fname,
                    },
                ))

        if not docs:
            # Presentation had slides but all were empty -- return a placeholder
            # so the upstream caller does not treat a zero-doc result as a crash.
            docs = [Document(
                page_content=f"[Empty presentation: {fname}]",
                metadata={"source": path},
            )]

        return docs

    @staticmethod
    def _load_csv(path: str) -> list[Document]:
        """Load a CSV file as multiple Documents — one per chunk of rows.

        Returns one Document per chunk of N data rows (default 10), each
        prepended with the header row.  This avoids returning a single giant
        Document that would be chunked arbitrarily by the text splitter and
        could cut across row boundaries, losing row context.

        Empty CSVs return a single placeholder document.
        """
        import csv as _csv

        with open(path, newline="", encoding="utf-8", errors="replace") as fh:
            reader = _csv.reader(fh)
            rows = list(reader)

        if not rows:
            return [Document(page_content="(empty CSV)", metadata={"source": path})]

        header = rows[0]
        header_line = " | ".join(header)
        data_rows = rows[1:]

        if not data_rows:
            content = f"[CSV file: {Path(path).name}]\n\n{header_line}"
            return [Document(page_content=content, metadata={"source": path})]

        chunk_size = 10
        docs: list[Document] = []
        for start in range(0, len(data_rows), chunk_size):
            batch = data_rows[start:start + chunk_size]
            body = "\n".join(" | ".join(r) for r in batch)
            content = (
                f"[CSV file: {Path(path).name} — rows {start + 1}–{start + len(batch)}]\n\n"
                f"{header_line}\n{body}"
            )
            docs.append(Document(page_content=content, metadata={"source": path}))
        return docs

    def copy_to_workspace(
        self,
        source_path: str,
        doc_type: DocumentType,
        workspace_paths: WorkspacePaths,
    ) -> str:
        """Copy a file to the appropriate classified folder in the workspace.

        If an existing workspace copy with the same content is found, that path
        is returned immediately without creating a new copy.  This prevents the
        same file from being re-indexed every time the user clicks Apply, which
        would accumulate duplicate content in the vector store.

        Returns the path of the workspace copy (new or reused).
        """
        import filecmp

        src = Path(source_path)
        if not src.exists():
            raise IngestError(f"Source file not found: {source_path}")

        dest_folder = workspace_paths.doc_folders.get(doc_type, workspace_paths.uploads)
        dest_folder.mkdir(parents=True, exist_ok=True)
        dest_path = dest_folder / src.name

        # If the primary destination already exists and has identical content,
        # reuse it — no new copy needed.  This is the common case when the user
        # clicks Apply multiple times without changing their file list.
        if dest_path.exists():
            try:
                if filecmp.cmp(str(src), str(dest_path), shallow=False):
                    return str(dest_path)
            except OSError:
                pass  # fall through to normal copy/rename logic

        # Primary destination exists but holds DIFFERENT content — find a free
        # suffixed name.  Also reuse any suffixed copy whose content matches.
        counter = 1
        candidate = dest_path
        while candidate.exists() and counter <= _MAX_COLLISION_COUNTER:
            stem = src.stem
            suffix = src.suffix
            candidate = dest_folder / f"{stem}_{counter}{suffix}"
            if candidate.exists():
                try:
                    if filecmp.cmp(str(src), str(candidate), shallow=False):
                        return str(candidate)
                except OSError:
                    pass
            counter += 1

        if counter > _MAX_COLLISION_COUNTER:
            raise IngestError(
                f"Could not find a unique filename for {src.name} in {dest_folder} "
                f"after {_MAX_COLLISION_COUNTER} attempts. "
                "The destination folder may contain an unexpected number of files."
            )

        shutil.copy2(str(src), str(candidate))
        return str(candidate)

    def split_documents(
        self,
        docs: list[Document],
        doc_type: DocumentType = DocumentType.other,
    ) -> list[Document]:
        """Split documents using a multi-stage type-specific pipeline."""
        try:
            return run_splitting_pipeline(docs, doc_type)
        except Exception as exc:  # noqa: BLE001
            raise IngestError(f"Failed to split documents: {exc}") from exc

    def load_and_split_classified(
        self,
        classified_files: dict[DocumentType, list[str]],
        workspace_paths: WorkspacePaths | None = None,
    ) -> list[Document]:
        """Load and split files by their document type classification.

        Individual file failures are logged and skipped so one bad file does
        not abort the entire indexing run.  If *every* file fails, an
        :class:`~uacragent.domain.errors.IngestError` is raised with a
        consolidated summary of what went wrong.

        Args:
            classified_files: Mapping of DocumentType to list of file paths
            workspace_paths: If provided, copy files to classified folders

        Returns:
            Combined list of all chunks with doc_type metadata
        """
        all_chunks: list[Document] = []
        _failed: list[tuple[str, str]] = []  # (filename, reason)

        # Track resolved paths across ALL doc types so the same physical file
        # is never loaded and embedded twice even when the user adds it to
        # multiple document-type buckets.
        _globally_seen: set[str] = set()

        for doc_type, paths in classified_files.items():
            # Deduplicate within this doc type (preserve order).
            unique_paths = list(dict.fromkeys(paths))
            if len(unique_paths) < len(paths):
                logger.warning(
                    "Duplicate file paths found for %s; deduplicating (%d → %d).",
                    doc_type, len(paths), len(unique_paths),
                )

            for path in unique_paths:
                # Cross-type duplicate check: compare resolved source paths so
                # the same file under two doc types is only indexed once (under
                # whichever bucket it appears in first).
                resolved = str(Path(path).resolve())
                if resolved in _globally_seen:
                    logger.warning(
                        "File '%s' was already processed under another document "
                        "type and will not be indexed again to avoid duplicate "
                        "content in the vector store.",
                        Path(path).name,
                    )
                    continue
                _globally_seen.add(resolved)

                fname = Path(path).name
                try:
                    # ── S5: pre-load size guard ────────────────────────────
                    # Check file size before reading to prevent OOM on very
                    # large files.  Stat failure is tolerated — the real error
                    # (missing file, permission denied) surfaces in load_single_file.
                    try:
                        file_size = Path(path).stat().st_size
                    except OSError:
                        file_size = 0

                    if file_size >= _LARGE_FILE_BLOCK_BYTES:
                        raise IngestError(
                            f"'{fname}' is {file_size / 1_048_576:.0f} MB — "
                            f"files over {_LARGE_FILE_BLOCK_BYTES // 1_048_576} MB "
                            "cannot be loaded (memory limit). "
                            "Split the file into smaller parts and re-upload."
                        )
                    if file_size >= _LARGE_FILE_WARN_BYTES:
                        logger.warning(
                            "Large file '%s' (%d MB) — loading may be slow.",
                            fname, file_size // 1_048_576,
                        )

                    # ── Optionally copy to workspace ──────────────────────
                    if workspace_paths:
                        path = self.copy_to_workspace(path, doc_type, workspace_paths)

                    # ── Load raw documents ────────────────────────────────
                    docs = self.load_single_file(path)

                    # ── B7: detect image-only / scanned PDFs ──────────────
                    # PyPDFLoader returns pages with empty page_content when
                    # the PDF has no embedded text layer (e.g. scanned pages).
                    # This produces zero usable chunks -- detect it early and
                    # give a clear, actionable error instead of the generic
                    # "No document chunks" message.
                    # PPTX files that contain only images are handled by the
                    # OCR path in _load_pptx() -- skip this check for them.
                    if Path(path).suffix.lower() == ".pdf" and docs:
                        if all(not d.page_content.strip() for d in docs):
                            raise IngestError(
                                f"'{fname}' appears to be a scanned or image-only PDF "
                                "— no text could be extracted from it. "
                                "Convert it to a searchable PDF using an OCR tool "
                                "(e.g. Adobe Acrobat, Tesseract) and re-upload."
                            )

                    # ── Split with type-specific pipeline ─────────────────
                    chunks = self.split_documents(docs, doc_type)
                    all_chunks.extend(chunks)

                except IngestError as exc:
                    _failed.append((fname, str(exc)))
                    logger.warning(
                        "Skipping '%s' (%s): %s",
                        fname, doc_type.value, exc,
                    )

        # If every file failed, surface a consolidated error so the caller
        # sees what went wrong rather than just "no chunks produced".
        if _failed and not all_chunks:
            summary = "; ".join(f"{n}: {r}" for n, r in _failed)
            raise IngestError(
                f"All {len(_failed)} file(s) failed to load. Details: {summary}"
            )

        return all_chunks

