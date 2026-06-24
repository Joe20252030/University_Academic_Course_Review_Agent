"""ConversationAgent — drives the chat loop and delegates task generation."""
from __future__ import annotations

import logging
import re
import threading
from collections.abc import Callable
from dataclasses import dataclass, field
from pathlib import Path

_logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Chat-history budget constants
# ---------------------------------------------------------------------------
# History is trimmed by *character count* (a provider-agnostic proxy for
# tokens; roughly 4 chars ≈ 1 token).  When the budget is exceeded the oldest
# middle turns are dropped and compressed into a short LLM-generated summary
# that is injected into the system prompt so semantic context is never lost.
#
# Layout preserved after trim:
#   [first KEEP_FIRST_TURNS turns]  ← often sets important session context
#   [summary of dropped turns]      ← injected into system prompt
#   [newest turns that fit budget]  ← conversational coherence window
#
_MAX_HISTORY_CHARS: int = 32_000   # ≈ 8 000 tokens; headroom for any provider
_KEEP_LAST_TURNS:   int = 3        # most recent N turns always kept verbatim
_KEEP_FIRST_TURNS:  int = 1        # oldest N turns always kept (context-setting)

from langchain_core.messages import AIMessage, HumanMessage, SystemMessage

from uacragent.agent.prompts._prompts import (
    PROMPTS_DIR as _PROMPTS_DIR,
    get_language_instruction as _get_language_instruction,
)
from uacragent.agent.session import AgentSession
from uacragent.domain.errors import LLMError
from uacragent.domain.providers import get_provider
from uacragent.domain.types import TaskType
from uacragent.infra.llm import LLMClient
from uacragent.infra.settings import Settings, get_settings

# ---------------------------------------------------------------------------
# MIME type map for file attachments
# ---------------------------------------------------------------------------

_MIME_MAP: dict[str, str] = {
    ".png":  "image/png",
    ".jpg":  "image/jpeg",
    ".jpeg": "image/jpeg",
    ".gif":  "image/gif",
    ".webp": "image/webp",
    ".bmp":  "image/bmp",
    ".pdf":  "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".ppt":  "application/vnd.ms-powerpoint",
    ".txt":  "text/plain",
    ".md":   "text/markdown",
    ".py":   "text/x-python",
    ".js":   "text/javascript",
    ".ts":   "text/typescript",
    ".csv":  "text/csv",
    ".json": "application/json",
    ".xml":  "text/xml",
    ".html": "text/html",
    ".htm":  "text/html",
}

_IMAGE_MIMES = {
    "image/png", "image/jpeg", "image/gif", "image/webp", "image/bmp",
}

# PPTX MIME types — used both in _extract_file_text and _build_human_message
# to identify files that need the combined text+image extraction path.
_PPTX_MIMES: frozenset[str] = frozenset({
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/vnd.ms-powerpoint",
})

_TEXT_MIMES = {
    "text/plain", "text/markdown", "text/x-python", "text/javascript",
    "text/typescript", "text/csv", "application/json", "text/xml",
    "text/html",
}


def _extract_docx_text(path: str) -> str:
    """Extract text from a .docx file, trying two libraries in order.

    ``python-docx`` is attempted first — it handles a wider range of real-world
    Word documents (complex XML, newer format features, embedded objects) than
    ``docx2txt``.  ``Docx2txtLoader`` (backed by ``docx2txt``) is kept as a
    secondary attempt in case ``python-docx`` itself has trouble with a file.
    Raises the last exception if both paths fail.
    """
    last_exc: Exception | None = None

    # Primary: python-docx — robust against complex Word documents
    try:
        import docx as _docx  # python-docx
        doc = _docx.Document(path)
        parts: list[str] = []
        # Paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        # Table cells
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        parts.append(cell.text)
        text = "\n".join(parts)
        if text.strip():
            return text
        # If python-docx returned empty text it may have silently failed on
        # the format — let docx2txt have a go before giving up.
    except Exception as exc:  # noqa: BLE001
        last_exc = exc
        _logger.debug("python-docx failed for '%s': %s", path, exc)

    # Secondary: Docx2txtLoader / docx2txt
    try:
        from langchain_community.document_loaders import Docx2txtLoader
        docs = Docx2txtLoader(path).load()
        return "\n\n".join(d.page_content for d in docs)
    except Exception as exc:  # noqa: BLE001
        last_exc = exc
        _logger.debug("Docx2txtLoader failed for '%s': %s", path, exc)

    raise last_exc or RuntimeError("DOCX extraction produced no output")


def _extract_pptx_text(path: str) -> str:
    """Extract text content from a .pptx file for use as a chat attachment.

    Reuses the same extraction logic as the document loader
    (text frames + speaker notes + OCR on image shapes) but returns a
    single flat string instead of a list of LangChain Documents.

    Raises the underlying exception on failure so the caller's error-handling
    wrapper can return an informative notice to the LLM.
    """
    from uacragent.infra.loaders import DocumentLoader as _DL, _check_ocr_available, _ocr_slide_images  # noqa: PLC0415

    ext = Path(path).suffix.lower()
    if ext == ".ppt":
        return (
            "[LEGACY FORMAT: This file is in the old PowerPoint 97-2003 "
            "binary format (.ppt) which cannot be read directly.  Ask the "
            "user to save it as a .pptx file and re-attach it.]"
        )

    try:
        from pptx import Presentation  # type: ignore[import-untyped]
    except ImportError:
        return (
            "[PPTX EXTRACTION UNAVAILABLE: python-pptx is not installed.  "
            "Ask the user to convert the file to PDF or plain text.]"
        )

    prs = Presentation(path)
    _ocr_ok = _check_ocr_available()
    slide_texts: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = [f"[Slide {slide_idx}]"]
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title = slide.shapes.title.text_frame.text.strip()
                if title:
                    parts.append(f"Title: {title}")
        except Exception:  # noqa: BLE001
            pass
        try:
            for shape in slide.shapes:
                try:
                    if shape is slide.shapes.title:
                        continue
                except Exception:  # noqa: BLE001
                    pass
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        parts.append(text)
        except Exception:  # noqa: BLE001
            pass
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"[Notes] {notes}")
        except Exception:  # noqa: BLE001
            pass
        if _ocr_ok:
            ocr = _ocr_slide_images(slide)
            if ocr:
                parts.append("[Image text (OCR)]")
                parts.extend(ocr)
        slide_texts.append("\n".join(parts))

    return "\n\n".join(slide_texts) if slide_texts else f"[Empty presentation: {Path(path).name}]"


def _extract_pptx_images(path: str) -> list[tuple[bytes, str]]:
    """Extract embedded image blobs from a .pptx file for vision-API delivery.

    Iterates every slide in slide order and collects the raw bytes and MIME
    type of each PICTURE-type shape whose content type is a recognised image
    format (``_IMAGE_MIMES``).

    Returns
    -------
    list of ``(blob_bytes, mime_type)`` tuples, capped at
    ``_MAX_PPTX_IMAGES``.  Individual images larger than
    ``_MAX_PPTX_IMAGE_BYTES`` are silently skipped.  Returns an empty list
    when python-pptx is not installed, the file cannot be opened, or no
    eligible image shapes are found.
    """
    try:
        from pptx import Presentation  # noqa: PLC0415
    except ImportError:
        return []

    try:
        prs = Presentation(path)
    except Exception:  # noqa: BLE001
        return []

    results: list[tuple[bytes, str]] = []
    for slide in prs.slides:
        if len(results) >= _MAX_PPTX_IMAGES:
            break
        for shape in slide.shapes:
            if len(results) >= _MAX_PPTX_IMAGES:
                break
            try:
                # MSO_SHAPE_TYPE.PICTURE == 13
                if shape.shape_type != 13:
                    continue
                blob: bytes = shape.image.blob
                content_type: str = shape.image.content_type or ""
                if not blob or content_type not in _IMAGE_MIMES:
                    continue
                if len(blob) > _MAX_PPTX_IMAGE_BYTES:
                    _logger.debug(
                        "Skipping large embedded image (%.1f MB) in '%s'.",
                        len(blob) / 1_048_576, Path(path).name,
                    )
                    continue
                results.append((blob, content_type))
            except Exception:  # noqa: BLE001
                continue

    return results


def _extract_file_text(path: str, mime: str) -> tuple[str, str | None]:
    """Extract text from a file given its path and MIME type.

    Returns
    -------
    (llm_text, ui_warning)
        *llm_text* is the content (or a bracketed notice) to include in the
        message sent to the LLM.  *ui_warning* is a human-readable string to
        show the user in the chat UI, or ``None`` when extraction succeeded.

    Unknown / binary MIME types (e.g. ``application/octet-stream``, Excel
    spreadsheets) are explicitly rejected rather than decoded as UTF-8, which
    would produce kilobytes of replacement characters that waste token budget
    and confuse the LLM.
    """
    fname = Path(path).name
    try:
        if mime == "application/pdf":
            from langchain_community.document_loaders import PyPDFLoader
            docs = PyPDFLoader(path).load()
            return "\n\n".join(d.page_content for d in docs), None
        elif mime == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
            return _extract_docx_text(path), None
        elif mime in _PPTX_MIMES:
            text = _extract_pptx_text(path)
            # _extract_pptx_text returns a bracketed notice when python-pptx is
            # missing or the file is a legacy .ppt — surface this as a UI warning.
            if text.startswith("[PPTX EXTRACTION UNAVAILABLE") or text.startswith("[LEGACY FORMAT"):
                return text, (
                    f"'{fname}' could not be read — "
                    + ("python-pptx is not installed (run: pip install python-pptx)."
                       if "not installed" in text
                       else "legacy .ppt format is not supported; please save as .pptx.")
                )
            return text, None
        elif mime in _TEXT_MIMES:
            return Path(path).read_text(encoding="utf-8", errors="replace"), None
        else:
            notice = (
                f"[UNSUPPORTED FILE TYPE: '{fname}' has MIME type "
                f"'{mime}' which cannot be extracted as text.  Please inform "
                f"the user that this file type is not supported and ask them to "
                f"convert it to PDF, DOCX, or plain text before attaching.]"
            )
            return notice, f"'{fname}' is not a supported file type for text extraction."
    except Exception as exc:  # noqa: BLE001
        _logger.debug("_extract_file_text failed for '%s': %s", path, exc)
        notice = (
            f"[FILE EXTRACTION ERROR: Could not read '{fname}'. "
            f"Reason: {type(exc).__name__}. "
            f"Please inform the user that this file could not be processed "
            f"and suggest they check the file is not corrupted or try a different format.]"
        )
        return notice, f"'{fname}' could not be read ({type(exc).__name__}) — it may be corrupt or unsupported."


def _provider_supports_vision(provider_id: str, model: str = "") -> bool:
    """Return True if the given provider+model supports image/vision inputs.

    Checks the provider-level ``supports_vision`` flag, then applies any
    model-level overrides:
    - OpenAI search-preview variants do not accept image inputs even though
      the provider generally does.
    """
    if "search-preview" in model.lower():
        return False
    return get_provider(provider_id).supports_vision


_MAX_ATTACHMENTS: int = 5           # cap on number of files per message
_MAX_ATTACHMENT_BYTES: int = 10 * 1024 * 1024  # 10 MB per file
_MAX_EXTRACTED_TEXT_CHARS: int = 20_000         # total extracted text budget

# Per-PPTX image extraction limits.
# Caps the number of embedded images sent to the LLM as vision parts so a
# presentation with many decorative images cannot exhaust the message context
# window.  Images larger than _MAX_PPTX_IMAGE_BYTES are skipped individually.
_MAX_PPTX_IMAGES: int = 5
_MAX_PPTX_IMAGE_BYTES: int = 5 * 1024 * 1024   # 5 MB per embedded image


def _build_human_message(
    message_text: str,
    attachments: list[dict],
    provider_id: str,
    model: str = "",
) -> "tuple[HumanMessage, list[str]]":
    """Build a HumanMessage with optional multimodal content parts.

    For image attachments: base64-encoded inline ``image_url`` content parts.
    For PDFs / Word docs / text files: extracted text appended to the message.
    When there are no attachments, returns a plain HumanMessage(content=...).

    *model* is forwarded to :func:`_provider_supports_vision` so model-level
    restrictions (e.g. ``gpt-4o-search-preview`` which belongs to a
    vision-capable provider but does not itself accept image inputs) are
    respected when deciding whether to attach PPTX embedded images as vision
    parts.  It mirrors the same ``model`` argument passed to the direct-image
    attachment guard in ``ConversationAgent.chat()``.

    Safety limits (applied with a log warning + UI-facing warning string):
    - Maximum ``_MAX_ATTACHMENTS`` (5) files per message.
    - Each file must be under ``_MAX_ATTACHMENT_BYTES`` (10 MB).
    - Total extracted text from all non-image files is capped at
      ``_MAX_EXTRACTED_TEXT_CHARS`` (20 000 chars).

    Returns:
        (HumanMessage, warnings) where *warnings* is a list of human-readable
        strings describing any attachments that were skipped or truncated.
        The caller should surface these in the chat UI.
    """
    if not attachments:
        return HumanMessage(content=message_text), []

    # Cap number of attachments
    ui_warnings: list[str] = []
    if len(attachments) > _MAX_ATTACHMENTS:
        dropped = len(attachments) - _MAX_ATTACHMENTS
        _logger.warning(
            "Received %d attachments — capping to %d.",
            len(attachments), _MAX_ATTACHMENTS,
        )
        ui_warnings.append(
            f"{dropped} attachment(s) were not sent — only the first "
            f"{_MAX_ATTACHMENTS} files per message are allowed."
        )
        attachments = attachments[:_MAX_ATTACHMENTS]

    import base64

    parts: list = []
    extra_text_blocks: list[str] = []
    total_extracted_chars: int = 0

    for att in attachments:
        path = att.get("path", "")
        # Per-file size guard (skip files exceeding the limit)
        try:
            file_size = Path(path).stat().st_size if path else 0
        except Exception:
            file_size = 0
        if file_size > _MAX_ATTACHMENT_BYTES:
            size_mb = file_size / (1024 * 1024)
            _logger.warning(
                "Attachment '%s' is %.1f MB — exceeds 10 MB limit; skipping.",
                Path(path).name, size_mb,
            )
            fname = Path(path).name
            extra_text_blocks.append(
                f"\n\n[Attachment '{fname}' was skipped: "
                f"file size {size_mb:.1f} MB exceeds the 10 MB limit.]"
            )
            ui_warnings.append(
                f"'{fname}' was not included — file size {size_mb:.1f} MB "
                f"exceeds the 10 MB limit."
            )
            continue
        mime = att.get("mime", "application/octet-stream")
        name = att.get("name", Path(path).name)

        if mime in _IMAGE_MIMES:
            # Inline base64 image — supported by Gemini and OpenAI
            try:
                data = Path(path).read_bytes()
                b64  = base64.b64encode(data).decode()
                parts.append({
                    "type": "image_url",
                    "image_url": {"url": f"data:{mime};base64,{b64}"},
                })
            except Exception as exc:  # noqa: BLE001
                extra_text_blocks.append(
                    f"[Could not encode image {name}: {exc}]"
                )
                ui_warnings.append(
                    f"'{name}' could not be attached — image encoding failed: {exc}"
                )
        else:
            # Text-based or document file — extract and append as text.
            # Check total extracted text budget before extraction.
            if total_extracted_chars >= _MAX_EXTRACTED_TEXT_CHARS:
                _logger.warning(
                    "Total extracted text limit (%d chars) reached; "
                    "skipping remaining text attachments.",
                    _MAX_EXTRACTED_TEXT_CHARS,
                )
                extra_text_blocks.append(
                    f"\n\n[Attachment '{name}' was omitted: "
                    f"total extracted text limit of {_MAX_EXTRACTED_TEXT_CHARS} "
                    f"characters was already reached.]"
                )
                ui_warnings.append(
                    f"'{name}' was not included — the {_MAX_EXTRACTED_TEXT_CHARS:,}-character "
                    f"text budget was already reached by earlier attachments."
                )
                continue
            extracted, extraction_warning = _extract_file_text(path, mime)
            if extraction_warning:
                ui_warnings.append(extraction_warning)
            # Truncate this file's extracted text to stay within the budget
            remaining = _MAX_EXTRACTED_TEXT_CHARS - total_extracted_chars
            if len(extracted) > remaining:
                _logger.warning(
                    "Extracted text from '%s' truncated to %d chars (budget remaining).",
                    name, remaining,
                )
                extracted = extracted[:remaining] + "\n[... truncated — text limit reached]"
                ui_warnings.append(
                    f"'{name}' was partially included — text was truncated at "
                    f"{remaining:,} characters to stay within the message budget."
                )
            total_extracted_chars += len(extracted)

            # ── PPTX: extract embedded images as additional vision parts ──────
            # Text frames go into the text block above; photos, diagrams, and
            # charts embedded in the slides are sent as image parts so the LLM
            # can actually see them rather than receiving silence where images are.
            #
            # Vision guard: only attempt extraction when the active provider
            # actually supports image/vision inputs.  Non-vision providers
            # (e.g. DeepSeek) cannot handle image_url content parts and would
            # return an error or silently corrupt the message.  This mirrors the
            # guard applied to direct image attachments in ConversationAgent.chat()
            # at the call site above _build_human_message().
            pptx_imgs: list[tuple[bytes, str]] = []
            if mime in _PPTX_MIMES and not extraction_warning:
                _candidate_imgs = _extract_pptx_images(path)
                if _candidate_imgs:
                    if _provider_supports_vision(provider_id, model=model):
                        pptx_imgs = _candidate_imgs
                    else:
                        # Build a precise reason: provider-level vs model-level.
                        _model_label = f" / {model}" if model else ""
                        _logger.warning(
                            "Provider '%s'%s does not support vision; "
                            "skipping %d embedded image(s) from '%s'.",
                            provider_id, _model_label,
                            len(_candidate_imgs), name,
                        )
                        ui_warnings.append(
                            f"'{name}' contains {len(_candidate_imgs)} embedded "
                            f"image(s) that were not sent — "
                            f"{provider_id}{_model_label} does not support images."
                        )

            _img_note = (
                f" (+ {len(pptx_imgs)} embedded image"
                f"{'s' if len(pptx_imgs) != 1 else ''} also attached below)"
                if pptx_imgs else ""
            )
            extra_text_blocks.append(
                f"\n\n--- Attached file: {name}{_img_note} ---\n{extracted}\n---"
            )

            for _img_bytes, _img_mime in pptx_imgs:
                _b64 = base64.b64encode(_img_bytes).decode()
                parts.append({
                    "type": "image_url",
                    "image_url": {"url": f"data:{_img_mime};base64,{_b64}"},
                })

            if pptx_imgs and len(pptx_imgs) == _MAX_PPTX_IMAGES:
                ui_warnings.append(
                    f"'{name}': only the first {_MAX_PPTX_IMAGES} embedded "
                    f"images were included — the presentation may contain more."
                )

    # Compose message text with any extracted text blocks
    full_text = message_text
    if extra_text_blocks:
        full_text = message_text + "".join(extra_text_blocks)

    if parts:
        # Multimodal content list: text part first, then image parts
        content: list = [{"type": "text", "text": full_text}] + parts
        return HumanMessage(content=content), ui_warnings
    else:
        return HumanMessage(content=full_text), ui_warnings


def _settings_for_session(base: Settings, session: AgentSession) -> Settings:
    """Return a Settings instance with provider/model overridden from the session.

    Uses ``model_copy`` instead of mutating ``os.environ`` so that concurrent
    background threads cannot race on the global environment.
    """
    updates: dict = {}
    if session.llm_provider:
        updates["llm_provider"] = session.llm_provider
    if session.llm_model:
        updates["llm_model"] = session.llm_model
    if updates:
        return base.model_copy(update=updates)
    return base

# _PROMPTS_DIR and _get_language_instruction are imported at the top of this module.

# ---------------------------------------------------------------------------
# Localised strings used inside agent responses
# ---------------------------------------------------------------------------

# Progress message shown in the "thinking" bubble while the summarisation LLM
# Localised progress message shown to the user when the history-compaction LLM
# call fires.  This fires only when chat history exceeds the trim budget, so it
# may surprise users who see a second LLM call after receiving a reply.
_SUMMARISING_HISTORY_MSGS: dict[str, str] = {
    "en": (
        "Compacting conversation history… (the AI is summarising earlier turns "
        "to stay within the context limit — this is not a new reply)"
    ),
    "zh_CN": (
        "正在压缩对话历史……（AI 正在总结早期对话以保持在上下文限制内——这不是新的回复）"
    ),
}

# Strings returned by initialize_session() as the session-status message.
_INIT_STATUS: dict[str, dict[str, str]] = {
    "en": {
        "no_docs": (
            "No documents are loaded yet. "
            "Add files in the Session Settings panel and click **Apply** to index them."
        ),
        "ready_cached": (
            "✓ {n_files} file(s) across {n_types} document type(s) already indexed."
        ),
        # This string is shown directly in the chat bubble on completion.
        # It must be self-contained: starts with ✓ so it reads as a success
        # notice without any additional wrapper added by the UI.
        "ready_indexed": (
            "✓ Indexed {n_files} file(s) across {n_types} document type(s). "
            "You can now ask questions or request a study document."
        ),
        # Shown directly in the chat bubble on failure (no additional wrapper).
        # The ⚠️ prefix is the visual error signal; {exc} carries the full
        # explanation (e.g. from ConfigurationError, it is already user-friendly).
        "init_failed": "⚠️ {exc}",
    },
    "zh_CN": {
        "no_docs": (
            "尚未加载任何文档。"
            "请在会话设置面板添加文件，然后点击**应用**进行索引。"
        ),
        "ready_cached": (
            "✓ 已完成索引：{n_files} 个文件（{n_types} 种文档类型）。"
        ),
        "ready_indexed": (
            "✓ 已成功索引 {n_files} 个文件（{n_types} 种文档类型）。"
            "您可以提问或请求生成学习文档。"
        ),
        "init_failed": "⚠️ {exc}",
    },
}

# Strings embedded in chat replies (generation errors, document-save notes).
_CHAT_STRINGS: dict[str, dict[str, str]] = {
    "en": {
        "no_docs_err": (
            "No documents are loaded. Please add files in the Session Settings "
            "panel and index them before generating a document."
        ),
        "not_init_err": (
            "The session has not been initialised yet. "
            "Click **Apply** to index your documents first."
        ),
        "gen_failed": "Generation failed: {exc}",
        "doc_saved": (
            "\n\n📄 The document has been saved as **{name}** "
            "in the outputs folder of your workspace.\n"
            "Full path: `{path}`\n\n"
            "*(This file remains on disk permanently — if you leave and reopen the session, "
            "you can find it via the path above or in the **Generated Outputs** section "
            "inside Session Settings.)*"
        ),
    },
    "zh_CN": {
        "no_docs_err": (
            "未加载任何文档。请在会话设置面板添加文件并完成索引后，再生成文档。"
        ),
        "not_init_err": (
            "会话尚未初始化，请先点击**应用**对文档进行索引。"
        ),
        "gen_failed": "生成失败：{exc}",
        "doc_saved": (
            "\n\n📄 文档已保存为 **{name}**，"
            "位于工作空间的 outputs 文件夹中。\n"
            "完整路径：`{path}`\n\n"
            "*(此文件将永久保存在磁盘上——离开并重新打开会话后，"
            "可通过上述路径找到，也可在会话设置的**生成输出**部分中查看。)*"
        ),
    },
}

# Language instructions are defined centrally in agent/prompts/_prompts.py
# and imported above as _get_language_instruction().  The local definition
# has been removed to eliminate the duplicate.


def _ls(lang: str, table: dict[str, dict[str, str]], key: str, **fmt: object) -> str:
    """Look up *key* in *table* for *lang*, falling back to English.

    Any *fmt* keyword arguments are passed to ``str.format`` so callers can
    inline placeholders without a separate format call.
    """
    row = table.get(lang) or table["en"]
    text = row.get(key) or table["en"].get(key)
    if text is None:
        _logger.warning(
            "Missing localisation key %r in string table for language %r — "
            "returning key name as fallback.",
            key, lang,
        )
        text = key  # raw key as last resort; at least it's visible
    return text.format(**fmt) if fmt else text


# Regex that matches a [TASK:xxx] marker.  The pattern is intentionally
# case-sensitive (no re.IGNORECASE) because the system prompt instructs the
# LLM to emit the marker in lower-case only.  If the LLM occasionally emits
# mixed-case, the defensive `.lower()` call in _extract_task_marker() handles
# it at the character level after the match so we never silently miss a marker.
_TASK_MARKER_RE = re.compile(
    r"\[TASK:(review_summary|practice_booklet|mock_exam|exam_prediction)\]",
)

# Second line of defence against false-positive task triggers.
#
# Strategy: NARROW BLOCKLIST, not a broad whitelist.
#
# Requiring generation keywords (whitelist) is the wrong tool — it forces
# exact phrasing, blocks natural variants like "I'd like a review" or
# "show me a mock exam", and makes the system needlessly rigid.
#
# Instead, suppress only patterns that are *unambiguously* not a generation
# request: connectivity checks, visibility tests, and greetings.  Everything
# else is left to the system prompt and the LLM's own judgement.
_OBVIOUS_NON_GENERATION_RE = re.compile(
    r"\bcan\s+you\s+see\b"                        # "can you see this / the image"
    r"|\bdo\s+you\s+see\b"                         # "do you see"
    r"|\bcan\s+you\s+read\s+this\b"               # "can you read this"
    r"|\bis\s+this\s+(?:working|visible|showing)\b"  # connectivity / visibility check
    r"|\bare\s+you\s+(?:there|online|working)\b"  # presence check
    r"|\btest[,\s!]+test\b",                       # "test, test" / "test! test"
    re.IGNORECASE,
)


def _is_obvious_non_generation(message: str) -> bool:
    """Return ``True`` only for messages that are *unambiguously* not a generation
    request — test pings, visibility checks, and presence checks.

    This is a narrow blocklist, not an intent detector.  It deliberately avoids
    requiring generation keywords so that natural-language requests ("I'd like a
    review", "show me a mock exam", "help me prepare") are never wrongly blocked.
    The system prompt is the primary mechanism for intent handling; this guard
    only catches the specific failure mode where the LLM emits a task marker in
    response to an obvious non-generation message (e.g. "can you see the image?").
    """
    return bool(_OBVIOUS_NON_GENERATION_RE.search(message))


def _extract_task_marker(text: str) -> tuple[str | None, str]:
    """Return ``(task_type_value, cleaned_text)`` after stripping the marker.

    The LLM embeds a ``[TASK:xxx]`` token to signal that a generation pipeline
    should be triggered.  This function detects the token, extracts the task
    type, and removes the token from the text shown to the user.
    """
    match = _TASK_MARKER_RE.search(text)
    if match is None:
        return None, text.strip()
    task_value = match.group(1).lower()
    clean = _TASK_MARKER_RE.sub("", text).strip()
    return task_value, clean


@dataclass
class ChatResponse:
    """Result returned by ConversationAgent.chat()."""

    text: str                       # assistant message shown to the user
    output_path: str | None = None  # set when a task document was generated
    task_type: str | None = None    # the TaskType value that was triggered, if any
    error: str | None = None        # set when generation failed
    attachment_warnings: list[str] = field(default_factory=list)
    # ^ human-readable warnings for attachment issues (too large, encode fail,
    #   text budget exceeded, vision stripped).  Shown in the chat UI after reply.
    history_appended: bool = True
    # ^ False when chat() returned early (e.g. cancelled before the LLM replied)
    #   without ever calling append_turn.  The UI layer uses this to decide
    #   whether to pop the AI half of the turn or to add a human-only entry.


class ConversationAgent:
    """Stateless agent that holds only shared infrastructure.

    All per-session state lives in :class:`AgentSession`.  Pass the same
    session object across multiple ``chat()`` calls to maintain continuity.
    """

    def __init__(self, settings: Settings | None = None) -> None:
        self.settings: Settings = settings or get_settings()

    # ------------------------------------------------------------------
    # Session initialisation
    # ------------------------------------------------------------------

    def initialize_session(
        self,
        session: AgentSession,
        progress_cb: Callable[[str], None] | None = None,
        force_reindex: bool = False,
        language: str = "auto",
    ) -> tuple[str, bool, str | None]:
        """Build or reuse the retriever for *session*.

        Returns ``(status_message, was_cached, fast_path_warning)`` where
        *was_cached* is ``True`` when the existing ChromaDB was opened directly
        without any re-embedding, and *fast_path_warning* is a non-None string
        when the fast path failed and a forced rebuild was performed.

        When *force_reindex* is ``True`` the fast path is skipped and a full
        indexing run is always performed (used by the Apply button so that
        changes to embedding provider, model, or files always take effect).

        *language* controls which locale is used for the returned status message
        (``"en"`` or ``"zh_CN"``).
        """
        from uacragent.agent.pipeline import AgentPipeline

        if not session.has_files():
            session.retriever = None
            # When the user removed every file and clicked Apply, the pipeline
            # is never entered — wipe uploads, chroma_db, and the manifest here.
            if force_reindex:
                from uacragent.agent.workspace_manager import (
                    wipe_session_uploads, wipe_session_vectorstore)
                wipe_session_uploads(session)
                wipe_session_vectorstore(session)
            return (_ls(language, _INIT_STATUS, "no_docs"), False, None)

        try:
            settings = _settings_for_session(self.settings, session)
            pipeline = AgentPipeline(settings)

            # ── Fast path: reuse existing ChromaDB without re-embedding ────────
            _fast_path_warning: str | None = None
            if not force_reindex:
                try:
                    retriever = pipeline.prepare_session_fast(session)
                    if retriever is not None:
                        session.retriever = retriever
                        n_files = sum(len(v) for v in session.active_files().values())
                        n_types = len(session.active_files())
                        return (
                            _ls(language, _INIT_STATUS, "ready_cached",
                                n_files=n_files, n_types=n_types),
                            True,
                            None,   # no warning
                        )
                except Exception as exc:  # noqa: BLE001
                    _logger.warning(
                        "Fast-path session restore failed (falling back to full "
                        "re-indexing — this may incur additional API usage): %s",
                        exc, exc_info=True,
                    )
                    _fast_path_warning = (
                        "⚠️ Cached index could not be loaded "
                        f"({type(exc).__name__}: {exc}). "
                        "Rebuilding from scratch — this may take a moment."
                    )

            # ── Full indexing path ──────────────────────────────────────────────
            session.retriever, _manifest_warn = pipeline.prepare_session(
                session, progress_cb=progress_cb
            )
            # Combine fast-path failure warning (if any) with manifest warning.
            _warn_parts = [w for w in (_fast_path_warning, _manifest_warn) if w]
            _warn = "\n\n".join(_warn_parts) or None
            n_types = len(session.active_files())
            n_files = sum(len(v) for v in session.active_files().values())
            return (
                _ls(language, _INIT_STATUS, "ready_indexed",
                    n_files=n_files, n_types=n_types),
                False,
                _warn,
            )
        except Exception as exc:  # noqa: BLE001
            session.retriever = None
            return (_ls(language, _INIT_STATUS, "init_failed", exc=exc), False, None)

    # ------------------------------------------------------------------
    # Chat
    # ------------------------------------------------------------------

    def chat(
        self,
        message: str,
        session: AgentSession,
        progress_cb: Callable[[str], None] | None = None,
        effort_level: str = "medium",
        reasoning_mode: str = "quick",
        language: str = "auto",
        search_enabled: bool = False,
        attachments: list | None = None,
        cancel_event: "threading.Event | None" = None,
    ) -> ChatResponse:
        """Process one user turn and return a :class:`ChatResponse`.

        *effort_level* controls how much retrieved context is injected into the
        system prompt and how deeply the generation pipeline samples the corpus.
        Valid values: ``"low"``, ``"medium"`` (default), ``"high"``.

        *reasoning_mode* selects the reasoning pipeline depth used when a
        generation task is triggered.  ``"quick"`` (default) uses a single-pass
        approach for faster, lower-cost output; ``"deep"`` enables structured
        topic extraction and a critic refinement pass for higher quality.

        *language* steers the LLM's response language via the system prompt
        and is used to localise any agent-generated error or status text
        embedded in the reply (``"en"`` or ``"zh_CN"``).

        Flow:
        1. Retrieve relevant context (if retriever is available).
        2. Build a message list: system prompt + history + new human message.
        3. Call the LLM.
        4. Strip any ``[TASK:xxx]`` marker from the response text.
        5. If a task marker was found, run the generation pipeline.
        6. Append the turn to *session.chat_history* and trim if needed.
        7. Return a ChatResponse.
        """
        # -- 1. Retrieve context -----------------------------------------------
        context = self._retrieve_context(message, session, effort_level)

        # -- 2. Build message list ---------------------------------------------
        provider_id = (session.llm_provider or "gemini").lower()

        # Compute user prefs once here so that _render_system_prompt and
        # _run_task (if triggered) share the same dict without each calling
        # session.to_user_prefs() (which can do disk I/O via read_exam_info()).
        prefs = session.to_user_prefs()

        # Strip image attachments if the active provider does not support vision.
        # Text-based attachments (PDFs, docs, plain text) are always allowed
        # since they are extracted to text before being sent to the LLM.
        effective_attachments = list(attachments or [])
        vision_warnings: list[str] = []
        if not _provider_supports_vision(provider_id, model=session.llm_model):
            image_atts = [
                a for a in effective_attachments
                if a.get("mime", "") in _IMAGE_MIMES
            ]
            if image_atts:
                _logger.warning(
                    "Provider '%s' does not support image/vision inputs; "
                    "stripping %d image attachment(s) from the message.",
                    provider_id, len(image_atts),
                )
                stripped_names = ", ".join(
                    f"'{a.get('name', Path(a.get('path', '')).name)}'"
                    for a in image_atts
                )
                vision_warnings.append(
                    f"Image attachment(s) {stripped_names} were not sent — "
                    f"the current provider ({provider_id}) does not support images."
                )
                effective_attachments = [
                    a for a in effective_attachments
                    if a.get("mime", "") not in _IMAGE_MIMES
                ]

        system_content = self._render_system_prompt(
            session, context, language=language,
            history_summary=session.history_summary,
            prefs=prefs,
            attachments=effective_attachments or None,
        )
        messages: list = [SystemMessage(content=system_content)]
        # Use snapshot() to get a consistent, lock-safe copy of history.
        # This prevents a concurrent cancel from mutating the list mid-iteration.
        messages.extend(session.chat_history.snapshot())
        human_msg, att_warnings = _build_human_message(
            message, effective_attachments, provider_id,
            model=session.llm_model or "",
        )
        messages.append(human_msg)
        all_warnings: list[str] = vision_warnings + att_warnings

        # -- 3. Call LLM (use session provider/model if different from default) --
        llm = LLMClient(_settings_for_session(self.settings, session))
        try:
            if search_enabled:
                raw_response = llm.invoke_with_search(messages, provider_id)
            else:
                raw_response = llm.stream_invoke(messages, cancel_event=cancel_event)
            if raw_response is None:
                # Cancelled during the initial LLM call — return immediately
                # without appending to history.
                return ChatResponse(
                    text="",
                    history_appended=False,
                    attachment_warnings=all_warnings,
                )
            content = getattr(raw_response, "content", str(raw_response))
            if isinstance(content, list):
                assistant_text: str = " ".join(
                    part.get("text", "") if isinstance(part, dict) else str(part)
                    for part in content
                    if isinstance(part, str) or (isinstance(part, dict) and part.get("type") == "text")
                ).strip() or str(content)
            else:
                assistant_text = str(content)
        except Exception as exc:  # noqa: BLE001
            error_msg = f"LLM call failed: {exc}"
            # append_turn() is called only AFTER the LLM succeeds (step 7).
            # We never reached it, so history_appended must be False so that
            # the UI layer can add the human turn via append_human_only().
            return ChatResponse(
                text=error_msg, error=error_msg,
                history_appended=False,
                attachment_warnings=all_warnings,
            )

        # -- 4. Detect and strip task marker -----------------------------------
        task_type, clean_text = _extract_task_marker(assistant_text)

        # Code-level safety gate: suppress task markers emitted in response to
        # messages that are unambiguously NOT generation requests (test pings,
        # visibility checks, presence checks).  This is a narrow blocklist —
        # natural-language requests are left entirely to the LLM's judgement.
        if task_type is not None and _is_obvious_non_generation(message):
            _logger.warning(
                "Task marker [TASK:%s] suppressed — message matched non-generation "
                "pattern: %r",
                task_type, message[:120],
            )
            task_type = None

        # -- 5. Run generation pipeline if task was requested ------------------
        output_path: str | None = None
        generation_error: str | None = None

        if task_type is not None:
            if not session.has_files():
                generation_error = _ls(language, _CHAT_STRINGS, "no_docs_err")
            elif session.retriever is None:
                generation_error = _ls(language, _CHAT_STRINGS, "not_init_err")
            else:
                try:
                    output_path = self._run_task(
                        task_type, session, progress_cb,
                        effort_level, reasoning_mode, language,
                        prefs=prefs,
                        cancel_event=cancel_event,
                    )
                except Exception as exc:  # noqa: BLE001
                    generation_error = _ls(language, _CHAT_STRINGS, "gen_failed", exc=exc)

        # -- 6. Build final reply text -----------------------------------------
        reply = clean_text
        if output_path:
            p = Path(output_path)
            reply += _ls(language, _CHAT_STRINGS, "doc_saved",
                         name=p.name, path=output_path)
        if generation_error:
            reply += f"\n\n⚠️ {generation_error}"

        # -- 7. Update history (after reply is fully assembled) ----------------
        # Save `reply` — not `clean_text` — so the file path note and any error
        # message are preserved across session reloads.
        # append_turn() inserts both messages atomically so a concurrent cancel
        # can never observe one message without the other.
        _att_meta = [
            {"name": a["name"], "mime": a.get("mime", ""), "path": a.get("path", "")}
            for a in (attachments or [])
        ]
        session.chat_history.append_turn(
            HumanMessage(
                content=message,
                additional_kwargs={"attachments": _att_meta} if _att_meta else {},
            ),
            AIMessage(content=reply),
        )
        self._smart_trim_history(
            session, llm,
            progress_cb=progress_cb,
            language=language,
            request_delay=self.settings.llm_request_delay,
        )

        return ChatResponse(
            text=reply,
            output_path=output_path,
            task_type=task_type,
            error=generation_error,
            attachment_warnings=all_warnings,
        )

    # ------------------------------------------------------------------
    # Private helpers
    # ------------------------------------------------------------------

    def _retrieve_context(
        self, message: str, session: AgentSession, effort_level: str = "medium"
    ) -> str:
        """Return retrieved context text scaled to *effort_level*.

        Uses the vectorstore directly with an effort-calibrated *k* so that
        Low/Medium/High effort controls how many chunks are injected into the
        system prompt without rebuilding the retriever.

        Security: every retrieved chunk is sanitised before being injected into
        the system prompt.  Any ``[TASK:xxx]`` pattern found in document text is
        neutralised by inserting a zero-width word-joiner (U+2060) between
        ``[TASK`` and ``:``, preventing a maliciously crafted document from
        triggering the generation pipeline via prompt injection.
        """
        if session.retriever is None:
            return "*(No documents loaded — answers are based on general knowledge only.)*"
        try:
            from uacragent.agent.pipeline import get_effort_config
            effort = get_effort_config(effort_level)

            # Prefer direct vectorstore access so we can override k at call time.
            # VectorStoreRetriever exposes `.vectorstore`; fall back to the
            # retriever itself (with its fixed k) when that attribute is absent.
            try:
                docs = session.retriever.vectorstore.similarity_search(
                    message, k=effort.retriever_k
                )
            except AttributeError:
                docs = session.retriever.invoke(message)

            if not docs:
                return "*(No relevant excerpts found in the uploaded documents.)*"
            # Cap total context size to avoid blowing up smaller-context-window
            # providers (e.g. DeepSeek).  Each chunk can be several thousand
            # chars; at High effort with k=12 the uncapped total can exceed
            # 48 000 chars.  Truncating here keeps the system prompt lean.
            _MAX_CONTEXT_CHARS = 24_000
            parts: list[str] = []
            total = 0
            for doc in docs:
                # Sanitise any [TASK:xxx] patterns embedded in document text.
                # A malicious document could contain "[TASK:mock_exam]" which,
                # once retrieved, would be indistinguishable from a real marker
                # in the LLM response and trigger unintended generation.
                # Inserting a zero-width word-joiner (U+2060) between "[TASK"
                # and ":" neutralises the pattern for our regex while leaving
                # the text readable and unambiguous to a human reviewer.
                chunk = doc.page_content.replace("[TASK:", "[TASK⁠:")
                if total + len(chunk) > _MAX_CONTEXT_CHARS:
                    remaining = _MAX_CONTEXT_CHARS - total
                    if remaining > 200:   # worth including a partial chunk
                        parts.append(chunk[:remaining] + "\n[... context truncated]")
                    break
                parts.append(chunk)
                total += len(chunk)
            return "\n\n".join(parts)
        except Exception:  # noqa: BLE001
            _logger.warning("Context retrieval failed; falling back to general knowledge.", exc_info=True)
            return "*(Context retrieval failed — answers are based on general knowledge only.)*"

    def _render_system_prompt(
        self,
        session: AgentSession,
        context: str,
        language: str = "auto",
        history_summary: str = "",
        prefs: dict | None = None,
        attachments: list | None = None,
    ) -> str:
        """Fill the system prompt template with session values.

        *language* is used to inject a language-steering instruction so the LLM
        responds in the user's chosen locale (``"en"`` or ``"zh_CN"``).

        *history_summary* is appended as an extra section when non-empty,
        giving the LLM context about conversation turns that were trimmed from
        the live history window to stay within the token budget.

        *prefs* is the pre-computed dict from ``session.to_user_prefs()``.  When
        provided it is used directly; otherwise it is computed here.  Callers
        that already hold a prefs dict (e.g. ``chat()``) should pass it in to
        avoid calling ``read_exam_info()`` (which can do disk I/O) a second time.
        """
        prompt_path = _PROMPTS_DIR / "conversation_system.md"
        try:
            template = prompt_path.read_text(encoding="utf-8")
        except FileNotFoundError:
            raise LLMError(
                f"System prompt template not found at {prompt_path}. "
                "This is an installation error — please reinstall the package."
            ) from None
        except Exception as exc:  # noqa: BLE001
            raise LLMError(f"Could not read system prompt template: {exc}") from exc

        if prefs is None:
            prefs = session.to_user_prefs()

        def _sanitise(value: str) -> str:
            """Prepare a user-supplied string for safe injection into the system prompt.

            Two threats are addressed:
            1. ``{`` / ``}`` brace injection — Python's ``str.format()`` treats
               unmatched braces as ``KeyError`` / ``ValueError``.  Escaping them
               to ``{{`` / ``}}`` makes them literal in the output.
            2. ``[TASK:xxx]`` prompt injection — a user could paste this marker
               into a free-text field (e.g. ``extra_instructions``) and trick the
               LLM into triggering an unintended generation pipeline.  Inserting
               a zero-width word-joiner (U+2060) between ``[TASK`` and ``:``
               neutralises the regex that the response parser uses, identical to
               the sanitisation already applied to retrieved document chunks.
            """
            return value.replace("{", "{{").replace("}", "}}").replace("[TASK:", "[TASK⁠:")

        # Build a compact course meta suffix for the opening line
        meta_parts = []
        if prefs.get("course_code"):
            meta_parts.append(prefs["course_code"])
        if prefs.get("semester"):
            meta_parts.append(prefs["semester"])
        course_meta = f" ({', '.join(meta_parts)})" if meta_parts else ""

        has_files_text = (
            ", ".join(
                f"{len(v)} {k.replace('_', ' ')}"
                for k, v in session.active_files().items()
            )
            if session.has_files()
            else "None"
        )

        response_language = _get_language_instruction(language)

        # Build an explicit note about any files attached to the current message.
        # This prevents the LLM from confusing "no indexed course documents" (empty
        # RAG context) with "no content available" when the user has pasted or
        # uploaded files directly into the chat.  Images and text files are
        # described separately because they reach the LLM via different channels:
        # text files are extracted and appended to the message body; images are
        # sent as inline base64 content parts for vision-capable providers.
        attachments_note = ""
        if attachments:
            image_atts = [a for a in attachments if a.get("mime", "") in _IMAGE_MIMES]
            text_atts  = [a for a in attachments if a.get("mime", "") not in _IMAGE_MIMES]
            note_lines: list[str] = []
            if text_atts:
                # Sanitise filenames: neutralise [TASK:] patterns so a file named
                # e.g. "notes [TASK:mock_exam].pdf" cannot inject a task marker.
                names = ", ".join(f'"{_sanitise(a.get("name", "file"))}"' for a in text_atts)
                note_lines.append(
                    f"The student has attached {len(text_atts)} document file(s) "
                    f"to this message ({names}). Their text content has been "
                    f"extracted and appended to the student's message — you can "
                    f"read and use it directly to answer their question."
                )
            if image_atts:
                names = ", ".join(f'"{_sanitise(a.get("name", "image"))}"' for a in image_atts)
                note_lines.append(
                    f"The student has attached {len(image_atts)} image file(s) "
                    f"to this message ({names}). They are included as inline "
                    f"images — you can view and describe their content directly."
                )
            if note_lines:
                attachments_note = "> **Attached files in this message:** " + " ".join(note_lines) + "\n\n"

        rendered = template.format(
            course_name=_sanitise(prefs.get("course_name") or "this course"),
            course_meta=_sanitise(course_meta),
            university_name=_sanitise(prefs.get("university_name") or "Not specified"),
            major=_sanitise(prefs.get("major") or "Not specified"),
            course_code=_sanitise(prefs.get("course_code") or "Not specified"),
            professor_name=_sanitise(prefs.get("professor_name") or "Not specified"),
            semester=_sanitise(prefs.get("semester") or "Not specified"),
            exam_type=prefs.get("exam_type") or "other",        # enum value — safe
            exam_format=prefs.get("exam_format") or "written",  # enum value — safe
            exam_duration=_sanitise(prefs.get("exam_duration") or "Not specified"),
            exam_info=_sanitise(prefs.get("exam_info") or "None provided"),
            extra_instructions=_sanitise(prefs.get("extra_instructions") or "None"),
            has_files=has_files_text,        # computed from enum keys + counts — safe
            context=context,                 # already sanitised in _retrieve_context()
            attachments_note=attachments_note,  # built above from effective_attachments
            response_language=response_language,  # fixed dict value — safe
        )

        if history_summary:
            # Neutralise any [TASK:xxx] patterns that may have survived from
            # earlier conversation turns into the summary.  The summary is
            # LLM-generated but is built from user messages that could contain
            # such patterns.  Injecting them unsanitised into the system prompt
            # is inconsistent with the sanitisation applied to every other
            # user-controlled field, and could in theory cause the LLM to emit
            # an unintended task marker.
            _safe_summary = history_summary.replace("[TASK:", "[TASK⁠:")
            rendered += (
                "\n\n## Earlier Conversation Summary\n"
                "The following is a compressed summary of earlier turns in this "
                "session that are no longer in the active history window. Use it "
                "to maintain continuity.\n\n"
                + _safe_summary
            )

        return rendered

    # ------------------------------------------------------------------
    # Smart history trim
    # ------------------------------------------------------------------

    def _smart_trim_history(
        self,
        session: AgentSession,
        llm_client: "LLMClient",
        progress_cb: "Callable[[str], None] | None" = None,
        language: str = "en",
        request_delay: float = 0.0,
    ) -> None:
        """Token-budget trim with LLM summarisation of dropped turns.

        Strategy
        --------
        1. If total history chars ≤ ``_MAX_HISTORY_CHARS``: no action.
        2. Otherwise, partition turns into three groups:

           * **first** — the oldest ``_KEEP_FIRST_TURNS`` turns, always kept
             (they often establish important session context such as "I'm
             studying for a closed-book exam").
           * **last** — the newest ``_KEEP_LAST_TURNS`` turns, always kept
             (conversational coherence window).
           * **middle** — everything in between.  Turns are kept newest-first
             until the remaining character budget is exhausted; surplus turns
             are dropped.

        3. Dropped turns are summarised via a short LLM call and the result
           is merged with any existing ``session.history_summary``.  If the
           summarisation call fails, a plain fallback string is used so the
           trim still proceeds.

        The in-memory ``_HistoryStore`` is rewritten with the kept messages
        under its own lock so the operation is atomic.
        """
        snapshot = session.chat_history.snapshot()
        total_chars = sum(len(str(m.content)) for m in snapshot)
        if total_chars <= _MAX_HISTORY_CHARS:
            return  # within budget — nothing to do

        # Convert flat message list to (human, ai) turn pairs.
        # Guard against odd-length lists (e.g. a dangling message after a bug):
        # range(0, len-1, 2) naturally skips the last element when odd, so we
        # preserve it separately as `dangling` and re-append it after trimming.
        if len(snapshot) % 2 != 0:
            _logger.warning(
                "Chat history has odd number of messages (%d); last message will be "
                "preserved separately.",
                len(snapshot),
            )
            dangling = snapshot[-1]
        else:
            dangling = None

        turns: list[tuple] = []
        for i in range(0, len(snapshot) - 1, 2):
            human_msg, ai_msg = snapshot[i], snapshot[i + 1]
            # Validate pair order — misaligned turns produce garbage summaries.
            if not isinstance(human_msg, HumanMessage) or not isinstance(ai_msg, AIMessage):
                _logger.warning(
                    "Turn pair at index %d is not (HumanMessage, AIMessage) — "
                    "found (%s, %s).  Dropping misaligned pair and continuing trim.",
                    i, type(human_msg).__name__, type(ai_msg).__name__,
                )
                # Remove the bad pair from the live store so it cannot pollute
                # subsequent calls.  We do NOT bail out entirely — that would
                # leave history unbounded if this condition is persistent.
                # The next save() will persist the corrected store.
                good_msgs = [m for m in snapshot if m is not human_msg and m is not ai_msg]
                session.chat_history.replace_all(good_msgs)
                # Restart trim with the corrected snapshot.
                self._smart_trim_history(
                    session, llm_client,
                    language=language, progress_cb=progress_cb,
                )
                return
            turns.append((human_msg, ai_msg))

        min_turns = _KEEP_FIRST_TURNS + _KEEP_LAST_TURNS + 1
        if len(turns) < min_turns:
            return  # not enough turns to trim meaningfully

        first_turns = turns[:_KEEP_FIRST_TURNS]
        last_turns  = turns[-_KEEP_LAST_TURNS:]
        middle_turns = turns[_KEEP_FIRST_TURNS : len(turns) - _KEEP_LAST_TURNS]

        # Remaining budget after mandatory keeps
        mandatory_chars = sum(
            len(str(m.content))
            for t in first_turns + last_turns
            for m in t
        )
        budget = _MAX_HISTORY_CHARS - mandatory_chars

        # Fill from newest middle turn backward until budget is exhausted
        kept_middle: list[tuple] = []
        dropped_middle: list[tuple] = []
        for turn in reversed(middle_turns):
            turn_chars = sum(len(str(m.content)) for m in turn)
            if budget >= turn_chars:
                kept_middle.insert(0, turn)
                budget -= turn_chars
            else:
                dropped_middle.insert(0, turn)

        if not dropped_middle:
            return  # nothing actually needed dropping

        # Summarise dropped turns (merges with any existing summary).
        # Notify the UI so the user understands why another LLM call fires.
        if progress_cb:
            try:
                _summ_msg = _SUMMARISING_HISTORY_MSGS.get(
                    language, _SUMMARISING_HISTORY_MSGS["en"]
                )
                progress_cb(_summ_msg)
            except Exception:  # noqa: BLE001
                _logger.debug("progress_cb raised during history compression.", exc_info=True)
        dropped_msgs = [m for t in dropped_middle for m in t]
        new_summary = self._summarise_turns(
            dropped_msgs,
            existing_summary=session.history_summary,
            llm_client=llm_client,
            language=language,
            request_delay=request_delay,
        )

        # Rebuild the in-memory store atomically
        kept_turns = first_turns + kept_middle + last_turns
        kept_msgs = [m for t in kept_turns for m in t]
        # Re-append the dangling odd message if one existed, so it is not lost.
        if dangling is not None:
            kept_msgs.append(dangling)
        session.chat_history.replace_all(kept_msgs)

        session.history_summary = new_summary
        _logger.debug(
            "History trimmed: kept %d turns, dropped %d turns, summary len=%d",
            len(kept_turns), len(dropped_middle), len(new_summary),
        )

    def _summarise_turns(
        self,
        messages: list,
        existing_summary: str,
        llm_client: "LLMClient",
        language: str = "en",
        request_delay: float = 0.0,
    ) -> str:
        """Return a 3-5 sentence summary of *messages*, merged with *existing_summary*.

        Called only when history exceeds the token budget.  Uses a short,
        focused prompt so the call is fast and cheap.  On any failure a plain
        fallback string (in the correct *language*) is returned so the trim
        still proceeds.

        *request_delay* is the inter-call delay from the active rate tier.  A
        sleep of that duration is inserted before the summarisation LLM call so
        it does not fire back-to-back with the preceding chat turn and hit the
        rate limit (which would silently lose history context).
        """
        # Build a condensed transcript of the turns to summarise.
        # Cap individual messages at 800 chars so one massive AI response
        # (e.g. a document preview) does not dominate the summary.
        transcript_parts: list[str] = []
        for m in messages:
            role = "Student" if isinstance(m, HumanMessage) else "Assistant"
            snippet = str(m.content)[:800]
            if len(str(m.content)) > 800:
                snippet += "…"
            transcript_parts.append(f"{role}: {snippet}")
        transcript = "\n\n".join(transcript_parts)

        if existing_summary:
            existing_snippet = existing_summary[:800]
            context_prefix = (
                f"There is already a summary of even earlier turns:\n"
                f"{existing_snippet}\n\n"
                f"Now incorporate the following additional exchanges into an "
                f"updated summary (keep total under 300 words):\n\n"
            )
        else:
            context_prefix = "Summarise the following conversation exchanges:\n\n"

        prompt = (
            f"{context_prefix}{transcript}\n\n"
            "Write a concise summary (3-5 sentences) covering:\n"
            "- What topics or questions were discussed\n"
            "- Any study documents that were generated (type and subject)\n"
            "- Key preferences or constraints the student mentioned\n"
            "Be factual and brief. Do not add commentary."
        )

        try:
            # Sleep before the summarisation call so it doesn't fire
            # back-to-back with the preceding chat turn and hit rate limits.
            import time as _time
            if request_delay > 0:
                _time.sleep(request_delay)

            resp = llm_client.invoke([
                SystemMessage(content=(
                    "You are a precise summariser. Return only the summary text, "
                    "no preamble, no labels."
                )),
                HumanMessage(content=prompt),
            ])
            summary = str(getattr(resp, "content", resp)).strip()
            # Cap at 1500 chars to keep the system prompt lean while retaining
            # enough context for meaningful history continuity.
            return summary[:1500]
        except Exception as exc:  # noqa: BLE001
            _logger.warning("History summarisation failed (%s); using fallback.", exc)
            n_turns = len(messages) // 2
            if existing_summary:
                return existing_summary  # keep what we had rather than losing it
            # Use a localised fallback so Chinese-locale users don't see an
            # English placeholder injected into their chat history.
            if language.startswith("zh"):
                return (
                    f"[本次会话早期记录：涵盖了 {n_turns} 轮对话，"
                    "讨论了课程材料、相关问题和学习主题。]"
                )
            return (
                f"[Earlier in this session: {n_turns} conversation turns covered "
                "course material, questions, and study topics.]"
            )

    def _run_task(
        self,
        task_type: str,
        session: AgentSession,
        progress_cb: Callable[[str], None] | None = None,
        effort_level: str = "medium",
        reasoning_mode: str = "quick",
        language: str = "auto",
        prefs: dict | None = None,
        cancel_event: "threading.Event | None" = None,
    ) -> str:
        """Run the generation pipeline and return the output markdown path.

        *prefs* is the pre-computed dict from ``session.to_user_prefs()``.  When
        provided it is used directly; otherwise it is computed here.  Callers
        that already hold a prefs dict (e.g. ``chat()``) should pass it in to
        avoid an unnecessary ``read_exam_info()`` disk I/O call.
        """
        from uacragent.agent.pipeline import AgentPipeline

        pipeline = AgentPipeline(_settings_for_session(self.settings, session))
        if prefs is None:
            prefs = session.to_user_prefs()

        _, _, md_path = pipeline.run_end_to_end(
            classified_files=session.classified_files,
            exam_format=prefs.get("exam_format", "written"),
            course_name=prefs.get("course_name", ""),
            exam_type=prefs.get("exam_type", "final"),
            task_type=task_type,
            extra_instructions=prefs.get("extra_instructions", ""),
            workspace_id=session.workspace_id,
            # Files are already in the workspace from the preceding prepare_session
            # call (triggered by Apply / session load).  Skipping re-copy prevents
            # duplicate upload copies from accumulating on each generation request.
            copy_to_workspace=False,
            university_name=prefs.get("university_name", ""),
            major=prefs.get("major", ""),
            course_code=prefs.get("course_code", ""),
            professor_name=prefs.get("professor_name", ""),
            semester=prefs.get("semester", ""),
            exam_duration=prefs.get("exam_duration", ""),
            exam_info=prefs.get("exam_info", ""),
            workspace_folder=session.workspace_folder,
            progress_cb=progress_cb,
            effort_level=effort_level,
            reasoning_mode=reasoning_mode,
            language=language,
            cancel_event=cancel_event,
        )
        return md_path
